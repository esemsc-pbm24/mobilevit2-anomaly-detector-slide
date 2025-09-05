import sys
import collections
# Compatibility shim for Python 3.10+ where ABCs moved to collections.abc
try:
    import collections.abc as _cabc
    for _name in ("Mapping", "MutableMapping", "Sequence", "Container"):
        if not hasattr(collections, _name) and hasattr(_cabc, _name):
            setattr(collections, _name, getattr(_cabc, _name))
except Exception:
    # Best-effort; if it fails, python-pptx import may still work on older Pythons
    pass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PASTEL_BLUE = RGBColor(0xA7, 0xD3, 0xF1)
PASTEL_GREEN = RGBColor(0xBF, 0xE7, 0xC6)
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)
LINE_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)

def set_run(run, text, size=18, bold=False, color=TEXT_DARK, name="Segoe UI"):
    run.text = text
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_badge(slide, text, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), text, size=18, bold=True)

def add_bullets(slide, left, top, width, height, sections):
    """
    sections: list of dicts:
      {"title": "Title", "subs": ["sub1", "sub2", ...]}
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    # spacing control
    for i, sec in enumerate(sections):
        # main bullet
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.level = 0
        p.space_after = Pt(4)
        set_run(p.add_run(), sec["title"], size=18, bold=True)
        # sub bullets
        for j, s in enumerate(sec.get("subs", [])):
            sp = tf.add_paragraph()
            sp.level = 1
            sp.space_after = Pt(2)
            set_run(sp.add_run(), s, size=14, bold=False)

def add_bottom_three_columns(slide, left, top, total_width, height, cols):
    """
    cols: list of dicts:
      {"badge": ("Label", color), "items": ["...", "..."]}
    """
    gap = Inches(0.25)
    col_w = (total_width - 2 * gap) / 3
    for i, col in enumerate(cols):
        x = left + i * (col_w + gap)
        # badge
        add_badge(slide, col["badge"][0], x, top, col_w, Inches(0.5), col["badge"][1])
        # bullets
        box_top = top + Inches(0.6)
        box = slide.shapes.add_textbox(x, box_top, col_w, height - Inches(0.6))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for j, item in enumerate(col["items"]):
            p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
            p.level = 0
            p.space_after = Pt(2)
            set_run(p.add_run(), f"• {item}", size=16, bold=False)

def add_divider(slide, left, top, width):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_LIGHT
    line.line.fill.background()

def main():
    prs = Presentation()
    # Set 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Slide background: light theme (white default)

    # Layout parameters
    margin = Inches(0.5)
    content_top = Inches(0.9)
    column_height = Inches(4.6)
    col_gap = Inches(0.4)
    col_width = (prs.slide_width - 2 * margin - col_gap) / 2

    left_x = margin
    right_x = margin + col_width + col_gap

    # Left badge and bullets
    add_badge(
        slide,
        "Advantages / Strengths",
        left_x,
        Inches(0.3),
        col_width,
        Inches(0.5),
        PASTEL_GREEN,
    )

    left_sections = [
        {
            "title": "High Recall Focus ✅",
            "subs": ["Critical for field deployment to reduce disease spread"],
        },
        {
            "title": "Modular Architecture 🔄",
            "subs": ["Autoencoder and classifier can work independently or together"],
        },
        {
            "title": "Ease of Data Acquisition 🌱",
            "subs": [
                "Autoencoder requires few or no diseased images",
                "Classifier trained only on images autoencoder fails to detect",
            ],
        },
    ]
    add_bullets(slide, left_x, content_top, col_width, column_height, left_sections)

    # Right badge and bullets
    add_badge(
        slide,
        "Limitations / Recommendations",
        right_x,
        Inches(0.3),
        col_width,
        Inches(0.5),
        PASTEL_BLUE,
    )

    right_sections = [
        {
            "title": "Autoencoder Limitations ⚠️",
            "subs": [
                "Weak decoder reduces standalone performance",
                "Normalization may worsen results",
            ],
        },
        {
            "title": "Improvement Opportunities 💡",
            "subs": [
                "Enhanced decoder (skip connections + attention mechanisms)",
                "Robust loss functions and tailored training strategies",
                "Ensemble with classifier for sparse proprietary data",
            ],
        },
    ]
    add_bullets(slide, right_x, content_top, col_width, column_height, right_sections)

    # Divider and Bottom section
    divider_y = content_top + column_height + Inches(0.2)
    add_divider(slide, margin, divider_y, prs.slide_width - 2 * margin)

    bottom_top = divider_y + Inches(0.2)
    bottom_height = prs.slide_height - bottom_top - Inches(0.4)

    bottom_cols = [
        {
            "badge": ("Dataset Observations 📸", PASTEL_BLUE),
            "items": [
                "Classifier handles varied conditions (day/night)",
                "More diseased samples → better performance",
            ],
        },
        {
            "badge": ("Model Choice 🏎️", PASTEL_GREEN),
            "items": [
                "Lightweight MobileViT2 outperforms larger models like EfficientNet",
            ],
        },
        {
            "badge": ("Takeaway ✔️", PASTEL_BLUE),
            "items": [
                "Lightweight anomaly detector + classifier is optimal for deployment",
                "Reduces data collection effort and maintains high recall",
            ],
        },
    ]
    add_bottom_three_columns(
        slide,
        margin,
        bottom_top,
        prs.slide_width - 2 * margin,
        bottom_height,
        bottom_cols,
    )

    out_name = "MobileViT2-Anomaly-Detector-Key-Insights.pptx"
    prs.save(out_name)
    print(f"Saved: {out_name}")


if __name__ == "__main__":
    main()